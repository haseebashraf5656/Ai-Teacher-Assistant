from fastapi import FastAPI, HTTPException, Depends, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel, EmailStr
from typing import Optional, List
import google.generativeai as genai
from datetime import datetime, timedelta
import jwt
import bcrypt
from sqlalchemy import create_engine, Column, Integer, String, DateTime, Text, ForeignKey, Float
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship
import os
import re
import json
import PyPDF2
from dotenv import load_dotenv

load_dotenv()  # Load environment variables from .env file

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PowerPoint file support disabled.")
import io

# Initialize FastAPI app
app = FastAPI(title="AI Teacher Assistant API")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure Gemini API
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
if not GEMINI_API_KEY:
    raise RuntimeError("GEMINI_API_KEY environment variable is not set. Please create a .env file.")
genai.configure(api_key=GEMINI_API_KEY)

# JWT Configuration
SECRET_KEY = os.getenv("SECRET_KEY", "change-this-secret-in-production")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 1440  # 24 hours

# Database Configuration
# For PostgreSQL (Production):
# DATABASE_URL = "postgresql://postgres:password@localhost:5432/teacher_assistant"
# engine = create_engine(DATABASE_URL)

# For SQLite (Development - Default):
DATABASE_URL = "sqlite:///./teacher_assistant.db"
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# Security
security = HTTPBearer()

# ==================== DATABASE MODELS ====================

class User(Base):
    __tablename__ = "users"
    
    id = Column(Integer, primary_key=True, index=True)
    email = Column(String, unique=True, index=True, nullable=False)
    full_name = Column(String, nullable=False)
    hashed_password = Column(String, nullable=False)
    created_at = Column(DateTime, default=datetime.utcnow)
    
    documents = relationship("Document", back_populates="owner")
    quizzes = relationship("Quiz", back_populates="owner")
    gradings = relationship("Grading", back_populates="owner")


class Document(Base):
    __tablename__ = "documents"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    filename = Column(String, nullable=False)
    file_type = Column(String, nullable=False)
    content = Column(Text, nullable=False)
    uploaded_at = Column(DateTime, default=datetime.utcnow)
    
    owner = relationship("User", back_populates="documents")
    quizzes = relationship("Quiz", back_populates="document")


class Quiz(Base):
    __tablename__ = "quizzes"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    document_id = Column(Integer, ForeignKey("documents.id"))
    difficulty = Column(String, nullable=False)
    selected_portion = Column(Text)
    questions = Column(Text, nullable=False)  # JSON string
    created_at = Column(DateTime, default=datetime.utcnow)
    
    owner = relationship("User", back_populates="quizzes")
    document = relationship("Document", back_populates="quizzes")


class Grading(Base):
    __tablename__ = "gradings"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    quiz_id = Column(Integer, ForeignKey("quizzes.id"))
    student_answer = Column(Text, nullable=False)
    marks_obtained = Column(Float, nullable=False)
    total_marks = Column(Float, nullable=False)
    feedback = Column(Text, nullable=False)
    graded_at = Column(DateTime, default=datetime.utcnow)
    
    owner = relationship("User", back_populates="gradings")


# Create all tables
Base.metadata.create_all(bind=engine)

# ==================== PYDANTIC MODELS ====================

class UserSignup(BaseModel):
    email: EmailStr
    full_name: str
    password: str


class UserLogin(BaseModel):
    email: EmailStr
    password: str


class Token(BaseModel):
    access_token: str
    token_type: str
    user_info: dict


class QuizGenerationRequest(BaseModel):
    document_id: int
    difficulty: str  # Easy, Medium, Hard
    num_questions: int = 3
    specific_requirements: Optional[str] = None


class QuizResponse(BaseModel):
    quiz_id: int
    questions: List[dict]
    message: str


class GradingRequest(BaseModel):
    quiz_id: int
    question_number: int
    student_answer: str
    total_marks: float


class GradingResponse(BaseModel):
    grading_id: int
    marks_obtained: float
    total_marks: float
    feedback: str
    percentage: float


# ==================== HELPER FUNCTIONS ====================

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')


def verify_password(plain_password: str, hashed_password: str) -> bool:
    return bcrypt.checkpw(plain_password.encode('utf-8'), hashed_password.encode('utf-8'))


def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt


def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security), db: Session = Depends(get_db)):
    token = credentials.credentials
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("sub")
        if email is None:
            raise HTTPException(status_code=401, detail="Invalid authentication credentials")
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token has expired")
    except jwt.JWTError:
        raise HTTPException(status_code=401, detail="Could not validate credentials")
    
    user = db.query(User).filter(User.email == email).first()
    if user is None:
        raise HTTPException(status_code=401, detail="User not found")
    return user


def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    pdf_file = io.BytesIO(file_content)
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def _extract_shape_text(shape) -> str:
    """Recursively extract text from a PPTX shape (including tables and groups)."""
    parts = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            child_text = _extract_shape_text(child)
            if child_text:
                parts.append(child_text)
        return "\n".join(parts)

    if shape.has_table:
        for row in shape.table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_cells.append(cell_text)
            if row_cells:
                parts.append(" | ".join(row_cells))
        return "\n".join(parts)

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = paragraph.text.strip()
            if paragraph_text:
                parts.append(paragraph_text)
        return "\n".join(parts)

    if hasattr(shape, "text") and shape.text.strip():
        parts.append(shape.text.strip())

    return "\n".join(parts)


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from modern PowerPoint (.pptx) files"""
    if not PPTX_AVAILABLE:
        raise HTTPException(
            status_code=400,
            detail="PowerPoint support not available. Run: pip install python-pptx"
        )
    pptx_file = io.BytesIO(file_content)
    presentation = Presentation(pptx_file)
    text = ""
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            shape_text = _extract_shape_text(shape)
            if shape_text.strip():
                slide_parts.append(shape_text.strip())
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"Notes: {notes_text}")
        if slide_parts:
            text += f"--- Slide {slide_num} ---\n" + "\n".join(slide_parts) + "\n"
    return text


def _extract_strings_from_ole(file_content: bytes) -> str:
    """Best-effort text extraction from legacy OLE Office files."""
    chunks = []
    seen = set()

    def add_chunk(value: str):
        cleaned = re.sub(r"\s+", " ", value).strip()
        if len(cleaned) < 4 or cleaned in seen:
            return
        alpha_count = sum(1 for char in cleaned if char.isalpha())
        if alpha_count < max(2, len(cleaned) // 4):
            return
        seen.add(cleaned)
        chunks.append(cleaned)

    try:
        decoded = file_content.decode("utf-16-le", errors="ignore")
        for match in re.findall(r"[A-Za-z0-9][A-Za-z0-9\s.,;:'\"()\-_/]{3,}", decoded):
            add_chunk(match)
    except Exception:
        pass

    for match in re.findall(rb"[\x20-\x7e]{4,}", file_content):
        add_chunk(match.decode("ascii", errors="ignore"))

    return "\n".join(chunks)


def extract_text_from_ppt_legacy(file_content: bytes) -> str:
    """Extract text from legacy .ppt (OLE) files."""
    temp_path = None
    try:
        import win32com.client
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as tmp:
            tmp.write(file_content)
            temp_path = tmp.name

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 0
        presentation = powerpoint.Presentations.Open(temp_path, WithWindow=False)
        text_parts = []
        for slide_index in range(1, presentation.Slides.Count + 1):
            slide = presentation.Slides(slide_index)
            slide_text = []
            for shape_index in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(shape_index)
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    slide_text.append(shape.TextFrame.TextRange.Text.strip())
            if slide_text:
                text_parts.append(f"--- Slide {slide_index} ---\n" + "\n".join(slide_text))
        presentation.Close()
        powerpoint.Quit()
        return "\n".join(text_parts)
    except ImportError:
        pass
    except Exception as exc:
        print(f"Legacy PPT COM extraction failed: {exc}")
    finally:
        if temp_path and os.path.exists(temp_path):
            os.unlink(temp_path)

    return _extract_strings_from_ole(file_content)


def _detect_file_type(file_content: bytes, content_type: str, filename: str) -> str:
    """Detect file type using magic bytes, filename extension, and MIME type."""
    fname = filename.lower() if filename else ""
    mime = (content_type or "").lower()

    if len(file_content) >= 4:
        if file_content[:4] == b"%PDF":
            return "pdf"
        if file_content[:4] == b"PK\x03\x04":
            return "pptx"
        if file_content[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
            return "ppt"

    if fname.endswith(".pdf"):
        return "pdf"
    if fname.endswith(".pptx"):
        return "pptx"
    if fname.endswith(".ppt"):
        return "ppt"
    if fname.endswith(".txt"):
        return "txt"

    if mime == "application/pdf":
        return "pdf"
    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return "pptx"
    if mime in [
        "application/vnd.ms-powerpoint",
        "application/powerpoint",
        "application/mspowerpoint",
        "application/x-mspowerpoint",
    ]:
        return "ppt"
    if mime in ["text/plain", "application/octet-stream"]:
        return "txt"

    return "unknown"


def extract_text_from_file(file_content: bytes, file_type: str, filename: str = "") -> str:
    """Extract text based on file type"""
    detected = _detect_file_type(file_content, file_type, filename)
    print(f"File: {filename}, MIME: {file_type}, Detected type: {detected}")

    if detected == "pdf":
        return extract_text_from_pdf(file_content)
    if detected == "pptx":
        return extract_text_from_pptx(file_content)
    if detected == "ppt":
        return extract_text_from_ppt_legacy(file_content)
    if detected == "txt":
        return file_content.decode("utf-8", errors="replace")
    raise HTTPException(
        status_code=400,
        detail=f"Unsupported file type. File: '{filename}', MIME: '{file_type}'. Supported formats: PDF, PPT, PPTX, TXT"
    )


def _normalize_document_content(content: str) -> str:
    """Clean extracted document text for AI processing."""
    cleaned = content.replace("\x0b", "\n").replace("\r\n", "\n").replace("\r", "\n")
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _parse_model_json(response_text: str) -> dict:
    """Parse JSON returned by Gemini, stripping optional markdown fences."""
    text = response_text.strip()
    if text.startswith("```json"):
        text = text[7:]
    elif text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    return json.loads(text.strip())


# ==================== API ENDPOINTS ====================

@app.get("/")
def root():
    return {
        "message": "AI Teacher Assistant API",
        "version": "1.0",
        "endpoints": {
            "signup": "/api/auth/signup",
            "login": "/api/auth/login",
            "upload_document": "/api/documents/upload",
            "generate_quiz": "/api/quiz/generate",
            "grade_answer": "/api/grading/grade"
        }
    }


@app.post("/api/auth/signup", response_model=Token)
def signup(user_data: UserSignup, db: Session = Depends(get_db)):
    """Register a new teacher"""
    # Check if user already exists
    existing_user = db.query(User).filter(User.email == user_data.email).first()
    if existing_user:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Create new user
    hashed_pwd = hash_password(user_data.password)
    new_user = User(
        email=user_data.email,
        full_name=user_data.full_name,
        hashed_password=hashed_pwd
    )
    
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    
    # Create access token
    access_token = create_access_token(data={"sub": new_user.email})
    
    return {
        "access_token": access_token,
        "token_type": "bearer",
        "user_info": {
            "id": new_user.id,
            "email": new_user.email,
            "full_name": new_user.full_name
        }
    }


@app.post("/api/auth/login", response_model=Token)
def login(user_data: UserLogin, db: Session = Depends(get_db)):
    """Login for existing teacher"""
    user = db.query(User).filter(User.email == user_data.email).first()
    
    if not user or not verify_password(user_data.password, user.hashed_password):
        raise HTTPException(status_code=401, detail="Incorrect email or password")
    
    # Create access token
    access_token = create_access_token(data={"sub": user.email})
    
    return {
        "access_token": access_token,
        "token_type": "bearer",
        "user_info": {
            "id": user.id,
            "email": user.email,
            "full_name": user.full_name
        }
    }


@app.post("/api/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Upload a document (PDF, PPT, or TXT)"""
    # Read file content
    file_content = await file.read()
    
    # Extract text from file
    try:
        text_content = extract_text_from_file(file_content, file.content_type, file.filename)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing file: {str(e)}")

    text_content = _normalize_document_content(text_content)

    if not text_content or len(text_content.strip()) < 20:
        raise HTTPException(
            status_code=400,
            detail="Could not extract enough text from the uploaded file. The file might be empty, image-only (scanned PDF/PPT), or corrupted. Try a text-based PDF or save the PowerPoint as .pptx."
        )

    print(f"Extracted {len(text_content)} characters from {file.filename}")
    print(f"Content preview (first 500 chars): {text_content[:500]}")
    
    # Save to database
    new_document = Document(
        user_id=current_user.id,
        filename=file.filename,
        file_type=file.content_type,
        content=text_content
    )
    
    db.add(new_document)
    db.commit()
    db.refresh(new_document)
    
    return {
        "message": "Document uploaded successfully",
        "document_id": new_document.id,
        "filename": new_document.filename,
        "content_length": len(text_content)
    }


@app.get("/api/documents/list")
def list_documents(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get list of all uploaded documents for the current user"""
    documents = db.query(Document).filter(Document.user_id == current_user.id).all()
    
    return {
        "documents": [
            {
                "id": doc.id,
                "filename": doc.filename,
                "file_type": doc.file_type,
                "uploaded_at": doc.uploaded_at.isoformat(),
                "content_preview": doc.content[:200] + "..." if len(doc.content) > 200 else doc.content
            }
            for doc in documents
        ]
    }


@app.post("/api/quiz/generate", response_model=QuizResponse)
def generate_quiz(
    request: QuizGenerationRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Generate quiz questions using Gemini AI"""
    # Get the document
    document = db.query(Document).filter(
        Document.id == request.document_id,
        Document.user_id == current_user.id
    ).first()
    
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")

    content = _normalize_document_content(document.content)
    if len(content) < 50:
        raise HTTPException(
            status_code=400,
            detail="This document does not contain enough readable text to generate a quiz. Please upload a text-based PDF, PPT, or PPTX file."
        )

    if request.difficulty not in ["Easy", "Medium", "Hard"]:
        raise HTTPException(status_code=400, detail="Difficulty must be Easy, Medium, or Hard")

    if request.num_questions < 1 or request.num_questions > 20:
        raise HTTPException(status_code=400, detail="Number of questions must be between 1 and 20")

    requirements_section = ""
    if request.specific_requirements and request.specific_requirements.strip():
        requirements_section = f"""
Teacher requirements (follow these while staying within the document content):
{request.specific_requirements.strip()}
"""

    content_to_send = content[:30000]
    print(f"Generating quiz: {request.num_questions} questions, {request.difficulty} difficulty")
    print(f"Document: {document.filename} ({len(content)} total chars, sending {len(content_to_send)} chars)")
    print(f"Document content preview: {content[:300]}...")

    prompt = f"""You are a teacher writing exam questions from the study material below.

<study_material filename="{document.filename}">
{content_to_send}
</study_material>

Create exactly {request.num_questions} {request.difficulty}-level questions using ONLY facts, terms, definitions, examples, and procedures found inside <study_material>.

Hard rules:
1. Every question and reference answer must be answerable using only the study material.
2. Do NOT ask about quizzes, exams, teachers, students, documents, files, or instructions for creating questions.
3. Do NOT invent topics that are not present in the study material.
4. Use specific terminology, names, numbers, and concepts that appear in the study material.
5. The reference answer must quote or paraphrase the study material accurately.
{requirements_section}
Difficulty guidance for {request.difficulty}:
- Easy: direct recall of facts, definitions, and terms from the material
- Medium: explain relationships, steps, or comparisons described in the material
- Hard: analyze or apply concepts from the material

Example of a good question if the material discusses Intel 8086 registers:
"What are the four general-purpose registers in the 8086/8088 architecture mentioned in the material?"

Example of a bad question (never do this):
"How many questions should this quiz contain?"

Return ONLY valid JSON in this exact format:
{{
    "questions": [
        {{
            "question_number": 1,
            "question_text": "Question based on the study material",
            "marks": 5,
            "reference_answer": "Expected answer based on the study material"
        }}
    ]
}}
"""

    try:
        model = genai.GenerativeModel("models/gemini-2.5-flash")
        response = model.generate_content(
            prompt,
            generation_config=genai.GenerationConfig(
                temperature=0.2,
                response_mime_type="application/json",
            ),
        )

        questions_data = _parse_model_json(response.text)
        if "questions" not in questions_data or not questions_data["questions"]:
            raise ValueError("Model response did not include any questions")

        new_quiz = Quiz(
            user_id=current_user.id,
            document_id=document.id,
            difficulty=request.difficulty,
            selected_portion=request.specific_requirements,
            questions=json.dumps(questions_data)
        )

        db.add(new_quiz)
        db.commit()
        db.refresh(new_quiz)

        return {
            "quiz_id": new_quiz.id,
            "questions": questions_data["questions"],
            "message": f"Successfully generated {len(questions_data['questions'])} questions at {request.difficulty} difficulty from {document.filename}"
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating quiz: {str(e)}")


@app.get("/api/quiz/list")
def list_quizzes(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get list of all quizzes for the current user"""
    quizzes = db.query(Quiz).filter(Quiz.user_id == current_user.id).all()
    
    import json
    return {
        "quizzes": [
            {
                "id": quiz.id,
                "document_id": quiz.document_id,
                "difficulty": quiz.difficulty,
                "created_at": quiz.created_at.isoformat(),
                "questions": json.loads(quiz.questions)
            }
            for quiz in quizzes
        ]
    }


@app.post("/api/grading/grade", response_model=GradingResponse)
def grade_answer(
    request: GradingRequest,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Grade a student's answer using Gemini AI"""
    # Get the quiz
    quiz = db.query(Quiz).filter(
        Quiz.id == request.quiz_id,
        Quiz.user_id == current_user.id
    ).first()
    
    if not quiz:
        raise HTTPException(status_code=404, detail="Quiz not found")
    
    # Parse quiz questions
    import json
    quiz_data = json.loads(quiz.questions)
    
    # Find the specific question
    question = None
    for q in quiz_data["questions"]:
        if q["question_number"] == request.question_number:
            question = q
            break
    
    if not question:
        raise HTTPException(status_code=404, detail="Question not found")
    
    # Create grading prompt
    prompt = f"""You are an expert teacher grading student answers.

Question: {question['question_text']}
Total Marks: {request.total_marks}
Reference Answer: {question.get('reference_answer', 'No reference answer provided')}

Student's Answer:
{request.student_answer}

Grade the student's answer and provide:
1. Marks obtained (out of {request.total_marks})
2. Detailed feedback explaining:
   - What was correct
   - What was incorrect or missing
   - Why marks were awarded or deducted

Be fair, constructive, and specific in your feedback.

Return your grading in the following JSON format:
{{
    "marks_obtained": 0.0,
    "feedback": "feedback here explaining the grading"
}}

Only return valid JSON, nothing else."""

    try:
        # Call Gemini API
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        response = model.generate_content(
            prompt,
            generation_config=genai.GenerationConfig(
                temperature=0.2,
                response_mime_type="application/json",
            ),
        )

        grading_data = _parse_model_json(response.text)
        
        marks_obtained = float(grading_data["marks_obtained"])
        feedback = grading_data["feedback"]
        
        # Ensure marks don't exceed total
        if marks_obtained > request.total_marks:
            marks_obtained = request.total_marks
        
        # Save grading to database
        new_grading = Grading(
            user_id=current_user.id,
            quiz_id=quiz.id,
            student_answer=request.student_answer,
            marks_obtained=marks_obtained,
            total_marks=request.total_marks,
            feedback=feedback
        )
        
        db.add(new_grading)
        db.commit()
        db.refresh(new_grading)
        
        percentage = (marks_obtained / request.total_marks) * 100
        
        return {
            "grading_id": new_grading.id,
            "marks_obtained": marks_obtained,
            "total_marks": request.total_marks,
            "feedback": feedback,
            "percentage": round(percentage, 2)
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error grading answer: {str(e)}")


@app.get("/api/grading/history")
def grading_history(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get grading history for the current user"""
    gradings = db.query(Grading).filter(Grading.user_id == current_user.id).all()
    
    return {
        "gradings": [
            {
                "id": grading.id,
                "quiz_id": grading.quiz_id,
                "marks_obtained": grading.marks_obtained,
                "total_marks": grading.total_marks,
                "percentage": round((grading.marks_obtained / grading.total_marks) * 100, 2),
                "feedback": grading.feedback,
                "graded_at": grading.graded_at.isoformat()
            }
            for grading in gradings
        ]
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)